import pandas as pd
import docx
import PyPDF2
from pptx import Presentation

def extract_text_txt(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except:
        return ""

def extract_text_docx(path):
    try:
        doc = docx.Document(path)
        return " ".join([p.text for p in doc.paragraphs])
    except:
        return ""

def extract_text_pdf(path):
    try:
        text = ""
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
        return text
    except:
        return ""

def extract_text_xlsx(path):
    try:
        df = pd.read_excel(path, dtype=str)
        return " ".join(df.fillna("").astype(str).values.flatten())
    except:
        return ""

def extract_text_csv(path):
    try:
        df = pd.read_csv(path, dtype=str)
        return " ".join(df.fillna("").astype(str).values.flatten())
    except:
        return ""

def extract_text_pptx(path):
    try:
        text = ""
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text
    except:
        return ""
